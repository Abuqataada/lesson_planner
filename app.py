import io
import json
import os
import re
import calendar
from contextlib import contextmanager
import unicodedata
import uuid
from datetime import datetime, timedelta, timezone
from pathlib import Path
from typing import Optional

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches, Pt
from fastapi import FastAPI, File, Form, Header, HTTPException, Query, Request, UploadFile
from fastapi.responses import HTMLResponse, Response
from fastapi.templating import Jinja2Templates
from dotenv import load_dotenv
import psycopg
from psycopg.rows import dict_row
from pydantic import BaseModel

try:
    from pypdf import PdfReader
except ImportError:
    PdfReader = None

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches as PptxInches
    from pptx.util import Pt as PptxPt
except ImportError:
    Presentation = None
    RGBColor = None
    MSO_SHAPE = None
    MSO_ANCHOR = None
    PP_ALIGN = None
    PptxInches = None
    PptxPt = None

try:
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT
    from reportlab.lib.pagesizes import landscape, letter
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import inch
    from reportlab.pdfbase.pdfmetrics import stringWidth
    from reportlab.lib.utils import simpleSplit
    from reportlab.pdfgen import canvas
    from reportlab.platypus import Paragraph, Table, TableStyle
except ImportError:
    colors = None
    TA_CENTER = None
    TA_LEFT = None
    TA_RIGHT = None
    landscape = None
    letter = None
    ParagraphStyle = None
    getSampleStyleSheet = None
    inch = None
    stringWidth = None
    simpleSplit = None
    canvas = None
    Paragraph = None
    Table = None
    TableStyle = None

from ai_analysis_service import AIAnalysisService

BASE_DIR = Path(__file__).resolve().parent
TEMPLATES_DIR = BASE_DIR / "templates"
load_dotenv(BASE_DIR / ".env")
DATABASE_URL = os.getenv(
    "DATABASE_URL"
)
ADMIN_PASSWORD = os.getenv("ADMIN_PASSWORD", "omtech-admin")
GENERATION_HISTORY_LIMIT = 10

PAYMENT_NAME = "OMTECH EI LIMITED"
PAYMENT_ACCOUNT = "7065894127"
PAYMENT_BANK = "MONIEPOINT"

SUBSCRIPTION_PLANS = [
    {
        "id": "monthly",
        "name": "Monthly Access",
        "price": "NGN 4,000",
        "billing": "per month",
        "badge": "Recommended",
        "description": "Best for schools and teachers who want flexible month-to-month access.",
    },
    {
        "id": "annual",
        "name": "Annual Access",
        "price": "NGN 40,000",
        "billing": "per annum",
        "badge": "Best Value",
        "description": "Best for active users who want uninterrupted access.",
    },
]

RENEWAL_NOTICE_DAYS = 7
GRACE_PERIOD_DAYS = 3

app = FastAPI(title="Lesson Planner")
templates = Jinja2Templates(directory=str(TEMPLATES_DIR))
ai_service = AIAnalysisService()


class LessonPlanRequest(BaseModel):
    class_name: str
    subject: str
    topic: str
    subscriber_key: Optional[str] = None
    output_format: str = "docx"


class SubscriptionRequest(BaseModel):
    full_name: str
    email: str
    phone: str
    plan_id: str = "monthly"
    payment_reference: str


def normalize_key(value: str) -> str:
    return (value or "").strip().lower()


def normalize_plan_id(value: str) -> str:
    key = normalize_key(value)
    aliases = {
        "termly": "monthly",
        "month": "monthly",
        "monthly": "monthly",
        "annual": "annual",
        "yearly": "annual",
        "annum": "annual",
    }
    return aliases.get(key, key or "monthly")


def iso_now() -> str:
    return datetime.utcnow().replace(microsecond=0).isoformat() + "Z"


def parse_iso_datetime(value: Optional[str]) -> Optional[datetime]:
    if not value:
        return None
    try:
        parsed = datetime.fromisoformat(value.replace("Z", "+00:00"))
        return parsed.astimezone(timezone.utc).replace(tzinfo=None)
    except Exception:
        return None


def format_iso_datetime(value: Optional[datetime]) -> Optional[str]:
    if not value:
        return None
    return value.replace(microsecond=0).isoformat() + "Z"


def add_months(base: datetime, months: int) -> datetime:
    base = base.replace(microsecond=0)
    month_index = base.month - 1 + months
    year = base.year + month_index // 12
    month = month_index % 12 + 1
    day = min(base.day, calendar.monthrange(year, month)[1])
    return base.replace(year=year, month=month, day=day)


def plan_cycle_months(plan_id: str) -> int:
    plan = normalize_plan_id(plan_id)
    return 12 if plan == "annual" else 1


def calculate_cycle_dates(plan_id: str, activated_at: Optional[str] = None) -> tuple[str, str, str]:
    activated_dt = parse_iso_datetime(activated_at) or datetime.utcnow().replace(microsecond=0)
    expires_dt = add_months(activated_dt, plan_cycle_months(plan_id))
    activated_iso = format_iso_datetime(activated_dt) or iso_now()
    expires_iso = format_iso_datetime(expires_dt) or iso_now()
    return activated_iso, expires_iso, expires_iso


def subscription_effective_state(record: Optional[dict]) -> dict:
    if not record:
        return {
            "status": "unregistered",
            "stored_status": None,
            "activated_at": None,
            "expires_at": None,
            "next_renewal_at": None,
            "grace_ends_at": None,
            "is_active": False,
            "in_grace_period": False,
            "days_remaining": None,
            "grace_days_remaining": None,
            "renewal_due_soon": False,
        }

    item = dict(record)
    stored_status = normalize_key(item.get("status"))
    plan_id = normalize_plan_id(item.get("plan_id", "monthly"))
    item["plan_id"] = plan_id

    activated_at = item.get("activated_at") or item.get("approved_at") or item.get("created_at")
    expires_at = item.get("expires_at") or item.get("next_renewal_at")
    activated_dt = parse_iso_datetime(activated_at) or parse_iso_datetime(item.get("created_at")) or datetime.utcnow().replace(microsecond=0)
    expires_dt = parse_iso_datetime(expires_at)

    days_remaining = None
    renewal_due_soon = False
    grace_ends_at = None
    grace_days_remaining = None
    in_grace_period = False

    if stored_status == "approved":
        if not expires_dt:
            activated_iso, expires_iso, renewal_iso = calculate_cycle_dates(plan_id, format_iso_datetime(activated_dt))
            activated_at = activated_iso
            expires_at = expires_iso
            next_renewal_at = renewal_iso
            expires_dt = parse_iso_datetime(expires_at)
        else:
            activated_at = format_iso_datetime(activated_dt)
            expires_at = format_iso_datetime(expires_dt)
            next_renewal_at = expires_at

        now = datetime.utcnow().replace(microsecond=0)
        grace_end_dt = expires_dt + timedelta(days=GRACE_PERIOD_DAYS) if expires_dt else None
        if expires_dt and now > grace_end_dt:
            effective_status = "expired"
            is_active = False
            grace_ends_at = format_iso_datetime(grace_end_dt)
        elif expires_dt and now > expires_dt:
            effective_status = "approved"
            is_active = True
            in_grace_period = True
            grace_ends_at = format_iso_datetime(grace_end_dt)
            delta = grace_end_dt - now
            grace_days_remaining = max(delta.days, 0)
            days_remaining = 0
        else:
            effective_status = "approved"
            is_active = True
            grace_ends_at = format_iso_datetime(grace_end_dt)
            if expires_dt:
                delta = expires_dt - now
                days_remaining = max(delta.days, 0)
                renewal_due_soon = delta <= timedelta(days=RENEWAL_NOTICE_DAYS)
    else:
        effective_status = stored_status or "pending"
        is_active = False
        activated_at = format_iso_datetime(parse_iso_datetime(activated_at))
        expires_at = format_iso_datetime(parse_iso_datetime(expires_at))
        next_renewal_at = expires_at

    item["stored_status"] = stored_status or "pending"
    item["status"] = effective_status
    item["activated_at"] = activated_at
    item["expires_at"] = expires_at
    item["next_renewal_at"] = next_renewal_at
    item["grace_ends_at"] = grace_ends_at
    item["is_active"] = is_active
    item["in_grace_period"] = in_grace_period
    item["days_remaining"] = days_remaining
    item["grace_days_remaining"] = grace_days_remaining
    item["renewal_due_soon"] = renewal_due_soon
    return item


@contextmanager
def get_db_connection():
    if not DATABASE_URL:
        raise RuntimeError("DATABASE_URL is not set. Configure the Neon/PostgreSQL connection in your environment.")
    if "YOUR_NEON_HOST" in DATABASE_URL or "YOUR_NEON_USER" in DATABASE_URL:
        raise RuntimeError(
            "DATABASE_URL still contains placeholder Neon values. Replace it in your .env or deployment environment."
        )
    conn = psycopg.connect(DATABASE_URL, row_factory=dict_row)
    try:
        yield conn
    finally:
        conn.close()


def init_db() -> None:
    with get_db_connection() as conn:
        conn.execute(
            """
            CREATE TABLE IF NOT EXISTS subscriptions (
                id TEXT PRIMARY KEY,
                subscriber_key TEXT NOT NULL UNIQUE,
                full_name TEXT NOT NULL,
                email TEXT NOT NULL,
                phone TEXT NOT NULL,
                plan_id TEXT NOT NULL,
                payment_reference TEXT NOT NULL,
                status TEXT NOT NULL,
                created_at TEXT NOT NULL,
                updated_at TEXT NOT NULL,
                activated_at TEXT,
                expires_at TEXT,
                next_renewal_at TEXT,
                approved_at TEXT,
                rejected_at TEXT
            )
            """
        )
        conn.execute("CREATE INDEX IF NOT EXISTS idx_subscriptions_status ON subscriptions(status)")
        conn.execute("ALTER TABLE subscriptions ADD COLUMN IF NOT EXISTS activated_at TEXT")
        conn.execute("ALTER TABLE subscriptions ADD COLUMN IF NOT EXISTS expires_at TEXT")
        conn.execute("ALTER TABLE subscriptions ADD COLUMN IF NOT EXISTS next_renewal_at TEXT")
        conn.execute(
            """
            CREATE TABLE IF NOT EXISTS generations (
                id TEXT PRIMARY KEY,
                subscriber_key TEXT NOT NULL,
                output_format TEXT NOT NULL,
                presentation_theme TEXT NOT NULL DEFAULT 'auto',
                subject TEXT NOT NULL,
                topic TEXT NOT NULL,
                filename TEXT NOT NULL,
                mime_type TEXT NOT NULL,
                file_data BYTEA NOT NULL,
                created_at TEXT NOT NULL
            )
            """
        )
        conn.execute("CREATE INDEX IF NOT EXISTS idx_generations_subscriber_created ON generations(subscriber_key, created_at DESC)")
        conn.execute("ALTER TABLE generations ADD COLUMN IF NOT EXISTS download_count INTEGER NOT NULL DEFAULT 0")
        conn.execute("ALTER TABLE generations ADD COLUMN IF NOT EXISTS last_download_at TEXT")
        conn.execute("ALTER TABLE generations ADD COLUMN IF NOT EXISTS presentation_theme TEXT NOT NULL DEFAULT 'auto'")
        conn.commit()


def load_subscriptions() -> list[dict]:
    with get_db_connection() as conn:
        rows = conn.execute(
            """
            SELECT id, subscriber_key, full_name, email, phone, plan_id, payment_reference,
                   status, created_at, updated_at, activated_at, expires_at, next_renewal_at, approved_at, rejected_at
            FROM subscriptions
            ORDER BY created_at DESC
            """
        ).fetchall()
    return [subscription_effective_state(dict(row)) for row in rows]


def get_generation_analytics() -> dict:
    with get_db_connection() as conn:
        totals = conn.execute(
            """
            SELECT
                COUNT(*) AS total_generations,
                COALESCE(SUM(download_count), 0) AS total_downloads
            FROM generations
            """
        ).fetchone()
        format_rows = conn.execute(
            """
            SELECT output_format, COUNT(*) AS count, COALESCE(SUM(download_count), 0) AS downloads
            FROM generations
            GROUP BY output_format
            ORDER BY count DESC
            """
        ).fetchall()
        subject_rows = conn.execute(
            """
            SELECT subject, COUNT(*) AS count, COALESCE(SUM(download_count), 0) AS downloads
            FROM generations
            GROUP BY subject
            ORDER BY count DESC, downloads DESC
            LIMIT 10
            """
        ).fetchall()
        theme_rows = conn.execute(
            """
            SELECT presentation_theme, COUNT(*) AS count
            FROM generations
            GROUP BY presentation_theme
            ORDER BY count DESC
            """
        ).fetchall()
        subscription_rows = conn.execute(
            """
            SELECT status, COUNT(*) AS count
            FROM subscriptions
            GROUP BY status
            """
        ).fetchall()
        recent_rows = conn.execute(
            """
            SELECT subject, topic, output_format, presentation_theme, filename, download_count, last_download_at, created_at
            FROM generations
            ORDER BY created_at DESC
            LIMIT 10
            """
        ).fetchall()
    subscription_counts = {row["status"]: row["count"] for row in subscription_rows}
    subscriptions = load_subscriptions()
    active_count = sum(1 for item in subscriptions if item.get("status") == "approved")
    expired_count = sum(1 for item in subscriptions if item.get("status") == "expired")
    renewal_due_count = sum(1 for item in subscriptions if item.get("status") == "approved" and item.get("renewal_due_soon"))
    grace_count = sum(1 for item in subscriptions if item.get("in_grace_period"))
    return {
        "summary": {
            "total_generations": int(totals["total_generations"] or 0),
            "total_downloads": int(totals["total_downloads"] or 0),
            "active_subscriptions": active_count,
            "expired_subscriptions": expired_count,
            "renewal_due_subscriptions": renewal_due_count,
            "grace_subscriptions": grace_count,
            "pending_subscriptions": int(subscription_counts.get("pending", 0)),
            "rejected_subscriptions": int(subscription_counts.get("rejected", 0)),
        },
        "by_format": [dict(row) for row in format_rows],
        "by_subject": [dict(row) for row in subject_rows],
        "by_theme": [dict(row) for row in theme_rows],
        "recent_generations": [dict(row) for row in recent_rows],
    }


def find_subscription(subscriber_key: str) -> Optional[dict]:
    key = normalize_key(subscriber_key)
    if not key:
        return None

    with get_db_connection() as conn:
        row = conn.execute(
            """
            SELECT id, subscriber_key, full_name, email, phone, plan_id, payment_reference,
                   status, created_at, updated_at, activated_at, expires_at, next_renewal_at, approved_at, rejected_at
            FROM subscriptions
            WHERE subscriber_key = %s OR lower(email) = %s OR lower(phone) = %s
            LIMIT 1
            """,
            (key, key, key),
        ).fetchone()
    return subscription_effective_state(dict(row)) if row else None


def upsert_subscription(payload: SubscriptionRequest) -> dict:
    now = iso_now()
    key = normalize_key(payload.email or payload.phone)
    existing = find_subscription(key)
    normalized_plan = normalize_plan_id(payload.plan_id)
    existing_active = bool(existing and existing.get("status") == "approved")

    record = {
        "id": existing["id"] if existing else str(uuid.uuid4()),
        "subscriber_key": key,
        "full_name": payload.full_name.strip(),
        "email": payload.email.strip(),
        "phone": payload.phone.strip(),
        "plan_id": normalized_plan,
        "payment_reference": payload.payment_reference.strip(),
        "status": "approved" if existing_active else "pending",
        "created_at": existing["created_at"] if existing else now,
        "updated_at": now,
        "activated_at": existing.get("activated_at") if existing_active else None,
        "expires_at": existing.get("expires_at") if existing_active else None,
        "next_renewal_at": existing.get("next_renewal_at") if existing_active else None,
        "approved_at": existing.get("approved_at") if existing_active else None,
        "rejected_at": None,
    }

    if existing_active and existing.get("status") == "approved" and existing.get("expires_at"):
        expires_dt = parse_iso_datetime(existing.get("expires_at"))
        if expires_dt and datetime.utcnow().replace(microsecond=0) > expires_dt:
            record["status"] = "pending"
            record["activated_at"] = None
            record["expires_at"] = None
            record["next_renewal_at"] = None
            record["approved_at"] = None

    with get_db_connection() as conn:
        conn.execute(
            """
            INSERT INTO subscriptions (
                id, subscriber_key, full_name, email, phone, plan_id, payment_reference,
                status, created_at, updated_at, activated_at, expires_at, next_renewal_at, approved_at, rejected_at
            ) VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s)
            ON CONFLICT(subscriber_key) DO UPDATE SET
                full_name=excluded.full_name,
                email=excluded.email,
                phone=excluded.phone,
                plan_id=excluded.plan_id,
                payment_reference=excluded.payment_reference,
                status=excluded.status,
                updated_at=excluded.updated_at,
                activated_at=excluded.activated_at,
                expires_at=excluded.expires_at,
                next_renewal_at=excluded.next_renewal_at,
                approved_at=excluded.approved_at,
                rejected_at=excluded.rejected_at
            """,
            (
                record["id"],
                record["subscriber_key"],
                record["full_name"],
                record["email"],
                record["phone"],
                record["plan_id"],
                record["payment_reference"],
                record["status"],
                record["created_at"],
                record["updated_at"],
                record["activated_at"],
                record["expires_at"],
                record["next_renewal_at"],
                record["approved_at"],
                record["rejected_at"],
            ),
        )
        conn.commit()

    return subscription_effective_state(record)


def set_subscription_status(subscription_id: str, status: str) -> dict:
    now = iso_now()
    with get_db_connection() as conn:
        row = conn.execute("SELECT * FROM subscriptions WHERE id = %s", (subscription_id,)).fetchone()
        if not row:
            raise HTTPException(status_code=404, detail="Subscription not found")
        item = dict(row)
        item["status"] = status
        item["updated_at"] = now
        if status == "approved":
            item["approved_at"] = now
            item["rejected_at"] = None
            item["activated_at"] = now
            activated_at, expires_at, next_renewal_at = calculate_cycle_dates(item.get("plan_id", "monthly"), now)
            item["activated_at"] = activated_at
            item["expires_at"] = expires_at
            item["next_renewal_at"] = next_renewal_at
        if status == "rejected":
            item["rejected_at"] = now
            item["approved_at"] = None
            item["activated_at"] = None
            item["expires_at"] = None
            item["next_renewal_at"] = None
        conn.execute(
            """
            UPDATE subscriptions
            SET status = %s, updated_at = %s, activated_at = %s, expires_at = %s, next_renewal_at = %s, approved_at = %s, rejected_at = %s
            WHERE id = %s
            """,
            (
                item["status"],
                item["updated_at"],
                item.get("activated_at"),
                item.get("expires_at"),
                item.get("next_renewal_at"),
                item.get("approved_at"),
                item.get("rejected_at"),
                subscription_id,
            ),
        )
        conn.commit()
        return subscription_effective_state(item)


def save_generation(
    subscriber_key: str,
    output_format: str,
    presentation_theme: str,
    subject: str,
    topic: str,
    filename: str,
    mime_type: str,
    file_bytes: bytes,
) -> dict:
    now = datetime.utcnow().isoformat(timespec="seconds") + "Z"
    record = {
        "id": str(uuid.uuid4()),
        "subscriber_key": normalize_key(subscriber_key),
        "output_format": normalize_output_format(output_format),
        "presentation_theme": normalize_key(presentation_theme or "auto"),
        "subject": subject.strip(),
        "topic": topic.strip(),
        "filename": filename,
        "mime_type": mime_type,
        "file_data": file_bytes,
        "created_at": now,
    }
    with get_db_connection() as conn:
        conn.execute(
            """
            INSERT INTO generations (
                id, subscriber_key, output_format, presentation_theme, subject, topic, filename, mime_type, file_data, created_at
            ) VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s, %s)
            """,
            (
                record["id"],
                record["subscriber_key"],
                record["output_format"],
                record["presentation_theme"],
                record["subject"],
                record["topic"],
                record["filename"],
                record["mime_type"],
                record["file_data"],
                record["created_at"],
            ),
        )
        conn.commit()
    return record


def list_generations(subscriber_key: str, limit: int = GENERATION_HISTORY_LIMIT) -> list[dict]:
    key = normalize_key(subscriber_key)
    if not key:
        return []
    with get_db_connection() as conn:
        rows = conn.execute(
            """
            SELECT id, subscriber_key, output_format, presentation_theme, subject, topic, filename, mime_type, created_at, download_count, last_download_at
            FROM generations
            WHERE subscriber_key = %s
            ORDER BY created_at DESC
            LIMIT %s
            """,
            (key, limit),
        ).fetchall()
    return [dict(row) for row in rows]


def get_generation(generation_id: str) -> Optional[dict]:
    with get_db_connection() as conn:
        row = conn.execute(
            """
            SELECT id, subscriber_key, output_format, presentation_theme, subject, topic, filename, mime_type, file_data, created_at, download_count, last_download_at
            FROM generations
            WHERE id = %s
            LIMIT 1
            """,
            (generation_id,),
        ).fetchone()
    return dict(row) if row else None


def increment_generation_download(generation_id: str) -> None:
    now = datetime.utcnow().isoformat(timespec="seconds") + "Z"
    with get_db_connection() as conn:
        conn.execute(
            """
            UPDATE generations
            SET download_count = COALESCE(download_count, 0) + 1,
                last_download_at = %s
            WHERE id = %s
            """,
            (now, generation_id),
        )
        conn.commit()


def admin_required(admin_password: str) -> None:
    if admin_password != ADMIN_PASSWORD:
        raise HTTPException(status_code=401, detail="Invalid admin password")


def clean_text(text: str) -> str:
    if not isinstance(text, str):
        text = str(text)
    replacements = {
        "\u2018": "'",
        "\u2019": "'",
        "\u201c": '"',
        "\u201d": '"',
        "\u2013": "-",
        "\u2014": "-",
        "\u2026": "...",
    }
    for bad, good in replacements.items():
        text = text.replace(bad, good)
    return unicodedata.normalize("NFKD", text).encode("ascii", "ignore").decode("ascii")


def sanitize_filename(text: str, max_len: int = 40) -> str:
    safe = "".join(c for c in text if c.isalnum() or c in (" ", "-", "_"))
    return safe.replace(" ", "_")[:max_len]


def normalize_output_format(value: str) -> str:
    fmt = normalize_key(value or "docx")
    if fmt in {"ppt", "pptx", "powerpoint"}:
        return "pptx"
    if fmt in {"pdf", "portable document format"}:
        return "pdf"
    return "docx"


def get_output_spec(output_format: str) -> dict[str, str]:
    fmt = normalize_output_format(output_format)
    if fmt == "pptx":
        return {
            "format": "pptx",
            "mime_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "extension": "pptx",
            "prefix": "Presentation",
        }
    if fmt == "pdf":
        return {
            "format": "pdf",
            "mime_type": "application/pdf",
            "extension": "pdf",
            "prefix": "Lesson_Export",
        }
    return {
        "format": "docx",
        "mime_type": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "extension": "docx",
        "prefix": "Lesson_Plan",
    }


def get_subject_theme(subject: str) -> dict[str, str]:
    subject_text = (subject or "").lower()
    if any(keyword in subject_text for keyword in ("math", "mathematics", "algebra", "geometry", "arithmetic", "statistics", "calculus")):
        return {
            "name": "Math",
            "primary": "#1d4ed8",
            "secondary": "#f59e0b",
            "accent": "#2563eb",
            "bg": "#eff6ff",
            "text": "#0f172a",
            "muted": "#475569",
        }
    if any(keyword in subject_text for keyword in ("science", "physics", "chemistry", "biology", "basic science", "integrated science")):
        return {
            "name": "Science",
            "primary": "#0f766e",
            "secondary": "#14b8a6",
            "accent": "#0ea5e9",
            "bg": "#ecfeff",
            "text": "#042f2e",
            "muted": "#0f766e",
        }
    if any(keyword in subject_text for keyword in ("english", "literature", "language", "grammar", "reading", "writing", "phonics")):
        return {
            "name": "English",
            "primary": "#166534",
            "secondary": "#84cc16",
            "accent": "#059669",
            "bg": "#f0fdf4",
            "text": "#052e16",
            "muted": "#166534",
        }
    return {
        "name": "General",
        "primary": "#2563eb",
        "secondary": "#14b8a6",
        "accent": "#2563eb",
        "bg": "#f8fafc",
        "text": "#0f172a",
        "muted": "#475569",
    }


PRESENTATION_THEME_PRESETS = [
    {
        "key": "auto",
        "name": "Auto Theme",
        "description": "Uses the best theme for the subject you enter.",
        "primary": "#2563eb",
        "secondary": "#14b8a6",
        "bg": "#eff6ff",
    },
    {
        "key": "science-lab",
        "name": "Science Lab",
        "description": "Teal and electric blue styling for science-related lessons.",
        "primary": "#0f766e",
        "secondary": "#14b8a6",
        "bg": "#ecfeff",
    },
    {
        "key": "math-grid",
        "name": "Math Grid",
        "description": "Crisp blue and amber visuals for mathematics and problem solving.",
        "primary": "#1d4ed8",
        "secondary": "#f59e0b",
        "bg": "#eff6ff",
    },
    {
        "key": "english-editorial",
        "name": "English Editorial",
        "description": "Clean green editorial styling for language and literature.",
        "primary": "#166534",
        "secondary": "#84cc16",
        "bg": "#f0fdf4",
    },
    {
        "key": "social-studies-heritage",
        "name": "Social Studies Heritage",
        "description": "Warm earth tones with a classroom-atlas feel.",
        "primary": "#7c2d12",
        "secondary": "#c2410c",
        "bg": "#fff7ed",
    },
    {
        "key": "ict-digital",
        "name": "ICT Digital",
        "description": "Modern indigo and cyan styling for computer and digital lessons.",
        "primary": "#312e81",
        "secondary": "#0ea5e9",
        "bg": "#eef2ff",
    },
]


PRESENTATION_THEME_LIBRARY = {
    "science-lab": {
        "name": "Science Lab",
        "primary": "#0f766e",
        "secondary": "#14b8a6",
        "accent": "#0ea5e9",
        "bg": "#ecfeff",
        "text": "#042f2e",
        "muted": "#0f766e",
        "accent_bg": "#dff9f6",
        "hero_note": "Scientific inquiry and practical exploration",
    },
    "math-grid": {
        "name": "Math Grid",
        "primary": "#1d4ed8",
        "secondary": "#f59e0b",
        "accent": "#2563eb",
        "bg": "#eff6ff",
        "text": "#0f172a",
        "muted": "#475569",
        "accent_bg": "#fff7ed",
        "hero_note": "Number sense, logic, and precision",
    },
    "english-editorial": {
        "name": "English Editorial",
        "primary": "#166534",
        "secondary": "#84cc16",
        "accent": "#059669",
        "bg": "#f0fdf4",
        "text": "#052e16",
        "muted": "#166534",
        "accent_bg": "#f7fee7",
        "hero_note": "Reading, writing, and language fluency",
    },
    "social-studies-heritage": {
        "name": "Social Studies Heritage",
        "primary": "#7c2d12",
        "secondary": "#c2410c",
        "accent": "#ea580c",
        "bg": "#fff7ed",
        "text": "#431407",
        "muted": "#9a3412",
        "accent_bg": "#ffedd5",
        "hero_note": "Culture, civic understanding, and community learning",
    },
    "ict-digital": {
        "name": "ICT Digital",
        "primary": "#312e81",
        "secondary": "#0ea5e9",
        "accent": "#6366f1",
        "bg": "#eef2ff",
        "text": "#111827",
        "muted": "#4338ca",
        "accent_bg": "#e0e7ff",
        "hero_note": "Technology skills, data, and digital confidence",
    },
    "default": {
        "name": "General",
        "primary": "#2563eb",
        "secondary": "#14b8a6",
        "accent": "#2563eb",
        "bg": "#f8fafc",
        "text": "#0f172a",
        "muted": "#475569",
        "accent_bg": "#eff6ff",
        "hero_note": "Balanced classroom presentation style",
    },
}


def resolve_presentation_theme(theme_key: str, subject: str) -> dict[str, str]:
    key = normalize_key(theme_key or "auto")
    subject_text = normalize_key(subject)
    if key in PRESENTATION_THEME_LIBRARY:
        return PRESENTATION_THEME_LIBRARY[key]
    if key == "auto":
        if any(keyword in subject_text for keyword in ("science", "physics", "chemistry", "biology", "basic science", "integrated science")):
            return PRESENTATION_THEME_LIBRARY["science-lab"]
        if any(keyword in subject_text for keyword in ("math", "mathematics", "algebra", "geometry", "arithmetic", "statistics", "calculus")):
            return PRESENTATION_THEME_LIBRARY["math-grid"]
        if any(keyword in subject_text for keyword in ("english", "literature", "language", "grammar", "reading", "writing", "phonics")):
            return PRESENTATION_THEME_LIBRARY["english-editorial"]
        if any(keyword in subject_text for keyword in ("social studies", "social studies", "civic", "government", "history", "geography")):
            return PRESENTATION_THEME_LIBRARY["social-studies-heritage"]
        if any(keyword in subject_text for keyword in ("ict", "computer", "computing", "digital", "technology", "information technology")):
            return PRESENTATION_THEME_LIBRARY["ict-digital"]
        return PRESENTATION_THEME_LIBRARY["default"]
    return PRESENTATION_THEME_LIBRARY["default"]


def extract_text_from_docx(file_bytes: bytes) -> str:
    document = Document(io.BytesIO(file_bytes))
    parts: list[str] = []
    for paragraph in document.paragraphs:
        if paragraph.text.strip():
            parts.append(paragraph.text.strip())
    for table in document.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                parts.append(row_text)
    return "\n".join(parts)


def extract_text_from_pdf(file_bytes: bytes) -> str:
    if PdfReader is None:
        raise HTTPException(status_code=500, detail="PDF extraction is unavailable. Install pypdf.")
    reader = PdfReader(io.BytesIO(file_bytes))
    parts: list[str] = []
    for page in reader.pages:
        text = (page.extract_text() or "").strip()
        if text:
            parts.append(text)
    return "\n".join(parts)


def extract_template_outline(upload: UploadFile) -> tuple[str, str]:
    file_bytes = upload.file.read()
    filename = (upload.filename or "").lower()
    if filename.endswith(".docx"):
        return extract_text_from_docx(file_bytes), "docx"
    if filename.endswith(".pdf"):
        return extract_text_from_pdf(file_bytes), "pdf"
    raise HTTPException(status_code=400, detail="Template must be a PDF or DOCX file.")


TEMPLATE_ALIASES = {
    "class": ["class", "class level", "class name", "grade"],
    "subject": ["subject"],
    "topic": ["topic"],
    "subtopic": ["subtopic", "sub topic"],
    "date": ["date"],
    "week": ["week"],
    "duration": ["duration", "period"],
    "age_group": ["age group", "student age group", "age range"],
    "resources": ["instructional resources", "resources", "teaching aids", "materials"],
    "learning objectives": ["learning objectives", "objectives"],
    "prior knowledge": ["prior knowledge", "previous knowledge"],
    "summarised learning note": ["summarised learning note", "learning note", "lesson note"],
    "warm-up": ["warm-up", "warm up", "starter", "introduction activity"],
    "teacher activities": ["teacher activities", "teacher's activities"],
    "student activities": ["student activities", "students' activities"],
    "assessment": ["assessment", "evaluation", "assessment/evaluation"],
    "plenary": ["plenary", "summary"],
    "homework": ["home-work", "homework", "assignment"],
    "flip ticket": ["flip ticket", "next topic", "preview"],
}


def derive_template_outline(template_text: str) -> dict[str, object]:
    raw_lines = [line.strip() for line in template_text.splitlines() if line.strip()]
    sections: list[str] = []
    labels: dict[str, str] = {}

    normalized_lines = [(line, normalize_key(line)) for line in raw_lines]
    for canonical, aliases in TEMPLATE_ALIASES.items():
        for original, normalized in normalized_lines:
            if any(alias in normalized for alias in aliases):
                if canonical not in sections:
                    sections.append(canonical)
                    labels[canonical] = original.rstrip(":") + ":"
                break

    return {
        "raw_lines": raw_lines[:80],
        "sections": sections,
        "labels": labels,
        "summary": " | ".join(sections) if sections else "default lesson structure",
        "prompt_text": "\n".join(raw_lines[:80]) if raw_lines else "default lesson structure",
    }


def landscape(doc: Document) -> None:
    section = doc.sections[0]
    section.page_width = Inches(11)
    section.page_height = Inches(8.5)
    section.top_margin = Inches(0.5)
    section.bottom_margin = Inches(0.5)
    section.left_margin = Inches(0.5)
    section.right_margin = Inches(0.5)


def create_lesson_plan_doc(
    plan_data: dict,
    teacher_name: str = "ISAH YUSUF",
    template_labels: Optional[dict[str, str]] = None,
    template_name: Optional[str] = None,
) -> bytes:
    template_labels = template_labels or {}
    doc = Document()
    landscape(doc)

    plan_data = {k: clean_text(v) if isinstance(v, str) else v for k, v in plan_data.items()}
    if isinstance(plan_data.get("learning_objectives"), dict):
        for key, value in list(plan_data["learning_objectives"].items()):
            if isinstance(value, str):
                plan_data["learning_objectives"][key] = clean_text(value)

    heading = doc.add_heading("Lesson Plan", level=1)
    heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
    if template_name:
        note = doc.add_paragraph(f"Template: {clean_text(template_name)}")
        note.alignment = WD_ALIGN_PARAGRAPH.CENTER
        note.runs[0].italic = True

    rows_data = [
        (template_labels.get("class", "Class:"), plan_data.get("class", "Year 7")),
        (template_labels.get("subject", "Subject:"), plan_data.get("subject", "Physics")),
        (template_labels.get("topic", "Topic:"), plan_data.get("topic", "Introduction to Physics")),
        (template_labels.get("subtopic", "Subtopic:"), plan_data.get("subtopic", plan_data.get("topic", ""))),
        (template_labels.get("date", "Date:"), plan_data.get("date", datetime.now().strftime("%d %B, %Y"))),
        (template_labels.get("week", "Week:"), plan_data.get("week", "Two")),
        (template_labels.get("duration", "Duration:"), plan_data.get("duration", "Forty Minutes")),
        (template_labels.get("age_group", "Student Age Group:"), plan_data.get("age_group", "11 - 12 Years")),
        (template_labels.get("resources", "Instructional Resources:"), ""),
    ]

    table_main = doc.add_table(rows=9, cols=3)
    table_main.style = "Table Grid"
    for i, (col1, col2) in enumerate(rows_data):
        table_main.cell(i, 0).text = clean_text(col1)
        if i == 8:
            cell2 = table_main.cell(i, 1)
            cell2.text = ""
            resources = plan_data.get("instructional_resources", [])
            if isinstance(resources, list):
                for res in resources:
                    p = cell2.add_paragraph(f"• {clean_text(res)}")
                    p.runs[0].font.size = Pt(11)
            else:
                p = cell2.add_paragraph(f"• {clean_text(resources)}")
                p.runs[0].font.size = Pt(11)
        else:
            table_main.cell(i, 1).text = clean_text(str(col2))

    col3_cell = table_main.cell(0, 2)
    for r in range(1, 9):
        col3_cell.merge(table_main.cell(r, 2))
    col3_cell.text = ""
    title = col3_cell.add_paragraph(template_labels.get("learning objectives", "Learning Objectives (Differentiated)"))
    title.runs[0].bold = True
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER

    lo = plan_data.get("learning_objectives", {})
    objective_blocks = [
        ("basic", "Basic Objective (for struggling learners):"),
        ("intermediate", "Intermediate Objective (for most students):"),
        ("advanced", "Advanced Objective (for high-achieving students):"),
    ]
    for key, label in objective_blocks:
        col3_cell.add_paragraph(template_labels.get(key, label)).runs[0].bold = True
        col3_cell.add_paragraph("By the end of the lesson, students will be able to:")
        bullet = col3_cell.add_paragraph(f"• {clean_text(lo.get(key, ''))}")
        bullet.paragraph_format.left_indent = Pt(36)

    doc.add_paragraph()
    table_dev = doc.add_table(rows=3, cols=6)
    table_dev.style = "Table Grid"

    prior_cell = table_dev.cell(0, 0)
    prior_cell.merge(table_dev.cell(1, 0))
    prior_cell.text = template_labels.get("prior knowledge", "Prior Knowledge")
    prior_cell.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER

    warmup_cell = table_dev.cell(0, 1)
    warmup_cell.merge(table_dev.cell(1, 1))
    warmup_cell.text = template_labels.get("warm-up", "Warm-up Activity")
    warmup_cell.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER

    note_cell = table_dev.cell(0, 2)
    note_cell.merge(table_dev.cell(1, 2))
    note_cell.text = template_labels.get("summarised learning note", "Summarised Learning Note")
    note_cell.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER

    learning_cell = table_dev.cell(0, 3)
    learning_cell.merge(table_dev.cell(0, 4))
    learning_cell.text = template_labels.get("learning activities", "Learning Activities")
    learning_cell.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER

    assess_cell = table_dev.cell(0, 5)
    assess_cell.merge(table_dev.cell(1, 5))
    assess_cell.text = template_labels.get("assessment", "Assessment/Evaluation")
    assess_cell.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER

    table_dev.cell(1, 3).text = template_labels.get("teacher activities", "TEACHER'S ACTIVITIES")
    table_dev.cell(1, 3).paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER
    table_dev.cell(1, 4).text = template_labels.get("student activities", "STUDENTS' ACTIVITIES")
    table_dev.cell(1, 4).paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER

    table_dev.cell(2, 0).text = clean_text(plan_data.get("prior_knowledge", "General knowledge from previous lessons."))
    table_dev.cell(2, 1).text = clean_text(plan_data.get("warmup_activity", "Engaging starter to capture interest."))
    table_dev.cell(2, 2).text = clean_text(plan_data.get("learning_note", "Core content with definitions and examples."))
    table_dev.cell(2, 3).text = clean_text(plan_data.get("teacher_activities", "Teacher's actions during the lesson."))
    table_dev.cell(2, 4).text = clean_text(plan_data.get("student_activities", "Group work, individual tasks, discussions."))
    table_dev.cell(2, 5).text = clean_text(plan_data.get("assessment", "Formative assessment methods."))

    doc.add_paragraph()
    table_plenary = doc.add_table(rows=3, cols=2)
    table_plenary.style = "Table Grid"
    plenary_rows = [
        (template_labels.get("plenary", "Plenary"), plan_data.get("plenary", "Summarise key points.")),
        (template_labels.get("homework", "Home-Work"), plan_data.get("homework", "Reinforcement task.")),
        (template_labels.get("flip ticket", "Flip Ticket (next Topic)"), plan_data.get("flip_ticket", "Preview of next lesson.")),
    ]
    for i, (label, text) in enumerate(plenary_rows):
        table_plenary.cell(i, 0).text = clean_text(label)
        table_plenary.cell(i, 1).text = clean_text(text)

    doc.add_paragraph()
    sig_row = doc.add_table(rows=1, cols=2)
    sig_row.style = "Table Grid"
    left_cell = sig_row.cell(0, 0)
    left_cell.paragraphs[0].text = f"Teacher's Name: {clean_text(teacher_name)}"
    left_cell.add_paragraph("Supervising Officer's Signature: ____________________")
    sig_row.cell(0, 1).text = (
        "Supervising officer's Comment:\n"
        "........................................................\n"
        "........................................................\n"
        "........................................................"
    )

    bio = io.BytesIO()
    doc.save(bio)
    bio.seek(0)
    return bio.getvalue()


def _require_reportlab() -> None:
    if canvas is None or colors is None or landscape is None or letter is None or simpleSplit is None:
        raise HTTPException(status_code=500, detail="PDF export is unavailable. Install reportlab.")


def _pdf_wrap_lines(text: str, font_name: str, font_size: int, max_width: float) -> list[str]:
    if not text:
        return [""]
    lines: list[str] = []
    for raw_line in str(text).splitlines() or [""]:
        wrapped = simpleSplit(raw_line, font_name, font_size, max_width)
        lines.extend(wrapped or [""])
    return lines or [""]


def _pdf_draw_wrapped(c, text: str, x: float, y: float, max_width: float, font_name: str = "Helvetica", font_size: int = 11, leading: int = 14, color_value=None) -> float:
    if color_value is None:
        color_value = colors.HexColor("#0f172a")
    c.setFillColor(color_value)
    c.setFont(font_name, font_size)
    line_y = y
    for line in _pdf_wrap_lines(clean_text(text), font_name, font_size, max_width):
        c.drawString(x, line_y, line)
        line_y -= leading
    return line_y


def _pdf_draw_bullets(c, items: list[str], x: float, y: float, max_width: float, font_name: str = "Helvetica", font_size: int = 11, leading: int = 14, bullet_color=None, text_color=None) -> float:
    bullet_color = bullet_color or colors.HexColor("#2563eb")
    text_color = text_color or colors.HexColor("#0f172a")
    line_y = y
    for item in items:
        c.setFillColor(bullet_color)
        c.setFont(font_name, font_size)
        c.drawString(x, line_y, "•")
        c.setFillColor(text_color)
        start_x = x + 14
        wrapped = _pdf_wrap_lines(clean_text(item), font_name, font_size, max_width - 14)
        for idx, line in enumerate(wrapped):
            c.drawString(start_x, line_y, line)
            line_y -= leading
            if idx < len(wrapped) - 1:
                line_y += leading
        line_y -= 4
    return line_y


def create_lesson_plan_pdf(
    plan_data: dict,
    teacher_name: str = "ISAH YUSUF",
    template_labels: Optional[dict[str, str]] = None,
    template_name: Optional[str] = None,
) -> bytes:
    _require_reportlab()
    template_labels = template_labels or {}
    theme = get_subject_theme(plan_data.get("subject", ""))
    buffer = io.BytesIO()
    page_width, page_height = landscape(letter)
    c = canvas.Canvas(buffer, pagesize=(page_width, page_height))

    def start_page(title: str, subtitle: str = ""):
        c.setFillColor(colors.HexColor(theme["bg"]))
        c.rect(0, 0, page_width, page_height, stroke=0, fill=1)
        c.setFillColor(colors.HexColor(theme["primary"]))
        c.rect(0, page_height - 42, page_width, 42, stroke=0, fill=1)
        c.setFillColor(colors.white)
        c.setFont("Helvetica-Bold", 16)
        c.drawString(24, page_height - 26, title)
        if subtitle:
            c.setFont("Helvetica", 10)
            c.drawRightString(page_width - 24, page_height - 26, subtitle)

    def card(x, y, w, h, fill="#ffffff", stroke="#dbeafe"):
        c.setFillColor(colors.HexColor(fill))
        c.setStrokeColor(colors.HexColor(stroke))
        c.roundRect(x, y, w, h, 14, stroke=1, fill=1)

    start_page("Lesson Plan", f"Theme: {theme['name']} | DOCX / PDF Ready")
    card(24, page_height - 192, page_width - 48, 110)
    c.setFillColor(colors.HexColor(theme["text"]))
    c.setFont("Helvetica-Bold", 22)
    c.drawString(40, page_height - 86, f"{clean_text(plan_data.get('subject', 'Subject'))} - {clean_text(plan_data.get('topic', 'Topic'))}")
    c.setFont("Helvetica", 11)
    c.setFillColor(colors.HexColor(theme["muted"]))
    c.drawString(40, page_height - 108, f"Class: {clean_text(plan_data.get('class', ''))}")
    c.drawString(40, page_height - 126, f"Date: {clean_text(plan_data.get('date', datetime.now().strftime('%d %B, %Y')))}")
    c.drawString(40, page_height - 144, f"Teacher: {clean_text(teacher_name)}")
    if template_name:
        c.drawString(40, page_height - 162, f"Template: {clean_text(template_name)}")

    left_y = page_height - 230
    card(24, 120, (page_width - 60) / 2, 300)
    card(page_width / 2 + 6, 120, (page_width - 60) / 2, 300)
    c.setFillColor(colors.HexColor(theme["primary"]))
    c.setFont("Helvetica-Bold", 14)
    c.drawString(40, 395, "Overview")
    c.drawString(page_width / 2 + 22, 395, "Teaching Flow")

    left_items = [
        f"{template_labels.get('prior knowledge', 'Prior Knowledge')}: {plan_data.get('prior_knowledge', '')}",
        f"{template_labels.get('warm-up', 'Warm-up Activity')}: {plan_data.get('warmup_activity', '')}",
        f"{template_labels.get('resources', 'Instructional Resources')}: " + ", ".join(_split_lines(plan_data.get("instructional_resources", []))),
        f"{template_labels.get('subtopic', 'Subtopic')}: {plan_data.get('subtopic', '')}",
    ]
    right_items = [
        f"{template_labels.get('summarised learning note', 'Summarised Learning Note')}: {plan_data.get('learning_note', '')}",
        f"{template_labels.get('teacher activities', 'Teacher Activities')}: {plan_data.get('teacher_activities', '')}",
        f"{template_labels.get('student activities', 'Student Activities')}: {plan_data.get('student_activities', '')}",
    ]
    _pdf_draw_bullets(c, left_items, 40, 370, (page_width - 92) / 2, font_size=10, leading=14, bullet_color=colors.HexColor(theme["secondary"]), text_color=colors.HexColor(theme["text"]))
    _pdf_draw_bullets(c, right_items, page_width / 2 + 22, 370, (page_width - 92) / 2, font_size=10, leading=14, bullet_color=colors.HexColor(theme["secondary"]), text_color=colors.HexColor(theme["text"]))

    c.showPage()
    start_page("Assessment and Review", "Teaching note")
    card(24, page_height - 250, page_width - 48, 170)
    c.setFillColor(colors.HexColor(theme["primary"]))
    c.setFont("Helvetica-Bold", 14)
    c.drawString(40, page_height - 86, "Assessment / Wrap-Up")
    assessment_lines = [
        f"{template_labels.get('assessment', 'Assessment')}: {plan_data.get('assessment', '')}",
        f"{template_labels.get('plenary', 'Plenary')}: {plan_data.get('plenary', '')}",
        f"{template_labels.get('homework', 'Homework')}: {plan_data.get('homework', '')}",
        f"{template_labels.get('flip ticket', 'Flip Ticket')}: {plan_data.get('flip_ticket', '')}",
    ]
    _pdf_draw_bullets(c, assessment_lines, 40, page_height - 118, page_width - 80, font_size=11, leading=16, bullet_color=colors.HexColor(theme["secondary"]), text_color=colors.HexColor(theme["text"]))

    card(24, 90, page_width - 48, 66, fill="#eff6ff", stroke="#bfdbfe")
    c.setFillColor(colors.HexColor(theme["primary"]))
    c.setFont("Helvetica-Bold", 11)
    c.drawString(40, 128, f"Prepared by {clean_text(teacher_name)}")
    c.setFont("Helvetica", 10)
    c.drawRightString(page_width - 40, 128, "Generated from the lesson planner")
    c.save()
    buffer.seek(0)
    return buffer.getvalue()


def create_presentation_pdf(
    plan_data: dict,
    teacher_name: str = "ISAH YUSUF",
    template_labels: Optional[dict[str, str]] = None,
    template_name: Optional[str] = None,
    theme_key: str = "auto",
) -> bytes:
    _require_reportlab()
    template_labels = template_labels or {}
    theme = resolve_presentation_theme(theme_key, plan_data.get("subject", ""))
    content = dict(plan_data or {})
    buffer = io.BytesIO()
    page_width, page_height = landscape(letter)
    c = canvas.Canvas(buffer, pagesize=(page_width, page_height))

    def slide_header(title: str, subtitle: str = ""):
        c.setFillColor(colors.HexColor(theme["bg"]))
        c.rect(0, 0, page_width, page_height, stroke=0, fill=1)
        c.setFillColor(colors.HexColor(theme["primary"]))
        c.rect(0, page_height - 46, page_width, 46, stroke=0, fill=1)
        c.setFillColor(colors.white)
        c.setFont("Helvetica-Bold", 16)
        c.drawString(24, page_height - 28, title)
        if subtitle:
            c.setFont("Helvetica", 10)
            c.drawRightString(page_width - 24, page_height - 28, subtitle)

    def slide_card(x, y, w, h, fill="#ffffff", stroke="#dbeafe"):
        c.setFillColor(colors.HexColor(fill))
        c.setStrokeColor(colors.HexColor(stroke))
        c.roundRect(x, y, w, h, 16, stroke=1, fill=1)

    slide_header(f"{clean_text(content.get('subject', 'Subject'))} Presentation", f"Theme: {theme['name']}")
    slide_card(24, 110, page_width - 48, page_height - 180, fill=theme["bg"], stroke="#bfdbfe")
    c.setFillColor(colors.HexColor(theme["text"]))
    c.setFont("Helvetica-Bold", 24)
    c.drawString(42, page_height - 92, f"{clean_text(content.get('topic', 'Topic'))}")
    c.setFont("Helvetica", 12)
    c.setFillColor(colors.HexColor(theme["muted"]))
    c.drawString(42, page_height - 114, f"Class: {clean_text(content.get('class', ''))} | Date: {clean_text(content.get('date', datetime.now().strftime('%d %B, %Y')))}")
    if template_name:
        c.drawString(42, page_height - 132, f"Template: {clean_text(template_name)}")
    _pdf_draw_wrapped(c, f"Teaching prompt: Open with the topic, then explain why it matters before moving to the meaning.", 42, page_height - 165, page_width - 84, font_size=11, leading=14, color_value=colors.HexColor(theme["primary"]))

    slide_card(42, 225, (page_width - 110) / 2, 220, fill="#ffffff", stroke="#dbeafe")
    slide_card(page_width / 2 + 14, 225, (page_width - 110) / 2, 220, fill="#ffffff", stroke="#dbeafe")
    c.setFillColor(colors.HexColor(theme["primary"]))
    c.setFont("Helvetica-Bold", 14)
    c.drawString(58, 424, "Overview")
    c.drawString(page_width / 2 + 30, 424, "Meaning")
    overview = [content.get("overview_line") or f"{content.get('topic', '')}: Definition, derivation of equations, and applications."]
    meaning = [content.get("meaning_text") or f"{content.get('topic', '')} is an important concept in {content.get('subject', '')}."]
    _pdf_draw_bullets(c, overview, 58, 398, (page_width - 130) / 2, font_size=10, leading=13, bullet_color=colors.HexColor(theme["secondary"]), text_color=colors.HexColor(theme["text"]))
    _pdf_draw_bullets(c, meaning, page_width / 2 + 30, 398, (page_width - 130) / 2, font_size=10, leading=13, bullet_color=colors.HexColor(theme["secondary"]), text_color=colors.HexColor(theme["text"]))

    c.showPage()
    slide_header(content.get("examples_heading") or "Examples", "Illustrations")
    slide_card(24, 118, page_width - 48, page_height - 188, fill="#ffffff", stroke="#dbeafe")
    c.setFillColor(colors.HexColor(theme["primary"]))
    c.setFont("Helvetica-Bold", 16)
    c.drawString(42, page_height - 90, "Examples")
    _pdf_draw_bullets(c, content.get("examples") or [f"An example related to {content.get('topic', '')}."], 42, page_height - 120, page_width - 84, font_size=11, leading=15, bullet_color=colors.HexColor(theme["secondary"]), text_color=colors.HexColor(theme["text"]))
    _pdf_draw_wrapped(c, f"Teaching prompt: Use the examples to link the topic to daily life and prior knowledge.", 42, 150, page_width - 84, font_size=11, leading=14, color_value=colors.HexColor(theme["primary"]))

    slide_card(24, 90, page_width - 48, 55, fill="#eff6ff", stroke="#bfdbfe")
    c.setFillColor(colors.HexColor(theme["primary"]))
    c.setFont("Helvetica-Bold", 11)
    c.drawString(42, 123, f"Prepared by {clean_text(teacher_name)}")
    c.setFont("Helvetica", 10)
    c.drawRightString(page_width - 42, 123, "Presentation export")
    c.showPage()

    slide_header(content.get("key_terms_heading") or "Terms Associated", "Key terms")
    slide_card(24, 118, page_width - 48, page_height - 188, fill="#ffffff", stroke="#dbeafe")
    c.setFillColor(colors.HexColor(theme["primary"]))
    c.setFont("Helvetica-Bold", 16)
    c.drawString(42, page_height - 90, content.get("key_terms_heading") or "Terms Associated")
    _pdf_draw_bullets(c, content.get("key_terms") or [f"Key term for {content.get('topic', '')}."], 42, page_height - 120, page_width - 84, font_size=11, leading=15, bullet_color=colors.HexColor(theme["secondary"]), text_color=colors.HexColor(theme["text"]))
    _pdf_draw_wrapped(c, f"Teaching prompt: Give learners time to copy the terms and ask a quick oral question.", 42, 150, page_width - 84, font_size=11, leading=14, color_value=colors.HexColor(theme["primary"]))

    c.showPage()
    slide_header(content.get("worked_examples_heading") or "Worked Examples", "Practice")
    slide_card(24, 118, page_width - 48, page_height - 188, fill="#ffffff", stroke="#dbeafe")
    c.setFillColor(colors.HexColor(theme["primary"]))
    c.setFont("Helvetica-Bold", 16)
    c.drawString(42, page_height - 90, content.get("worked_examples_heading") or "Worked Examples")
    _pdf_draw_bullets(c, content.get("worked_examples") or [f"Worked example on {content.get('topic', '')}."], 42, page_height - 120, page_width - 84, font_size=11, leading=15, bullet_color=colors.HexColor(theme["secondary"]), text_color=colors.HexColor(theme["text"]))
    _pdf_draw_wrapped(c, f"Teaching prompt: Walk through each worked example slowly and allow learners to solve the next one.", 42, 150, page_width - 84, font_size=11, leading=14, color_value=colors.HexColor(theme["primary"]))

    c.showPage()
    slide_header(content.get("classwork_heading") or "CLASSWORK", "In class exercise")
    slide_card(24, 118, page_width - 48, page_height - 188, fill="#ffffff", stroke="#dbeafe")
    c.setFillColor(colors.HexColor(theme["primary"]))
    c.setFont("Helvetica-Bold", 16)
    c.drawString(42, page_height - 90, content.get("classwork_heading") or "CLASSWORK")
    _pdf_draw_bullets(c, content.get("classwork") or [f"1. Define {content.get('topic', '')}."], 42, page_height - 120, page_width - 84, font_size=11, leading=15, bullet_color=colors.HexColor(theme["secondary"]), text_color=colors.HexColor(theme["text"]))
    _pdf_draw_wrapped(c, f"Teaching prompt: Pause for responses, then review answers with the class.", 42, 150, page_width - 84, font_size=11, leading=14, color_value=colors.HexColor(theme["primary"]))

    c.showPage()
    slide_header(content.get("weekend_assignment_heading") or "WEEKEND ASSIGNMENT", "Home practice")
    slide_card(24, 118, page_width - 48, page_height - 188, fill="#ffffff", stroke="#dbeafe")
    c.setFillColor(colors.HexColor(theme["primary"]))
    c.setFont("Helvetica-Bold", 16)
    c.drawString(42, page_height - 90, content.get("weekend_assignment_heading") or "WEEKEND ASSIGNMENT")
    _pdf_draw_bullets(c, content.get("weekend_assignment") or [f"1. Study {content.get('topic', '')} and answer related questions."], 42, page_height - 120, page_width - 84, font_size=11, leading=15, bullet_color=colors.HexColor(theme["secondary"]), text_color=colors.HexColor(theme["text"]))
    _pdf_draw_wrapped(c, f"Teaching prompt: Use this as the closing reinforcement task before the next lesson.", 42, 150, page_width - 84, font_size=11, leading=14, color_value=colors.HexColor(theme["primary"]))

    c.showPage()
    slide_header("THANK YOU", theme["name"])
    slide_card(24, 118, page_width - 48, page_height - 188, fill=theme["accent_bg"], stroke="#bfdbfe")
    c.setFillColor(colors.HexColor(theme["primary"]))
    c.setFont("Helvetica-Bold", 22)
    c.drawCentredString(page_width / 2, page_height - 95, "THANK YOU")
    c.setFont("Helvetica", 12)
    c.setFillColor(colors.HexColor(theme["text"]))
    c.drawCentredString(page_width / 2, page_height - 122, theme["hero_note"])
    c.drawCentredString(page_width / 2, page_height - 146, f"Prepared by {clean_text(teacher_name)}")
    c.drawCentredString(page_width / 2, page_height - 170, "Presentation export")

    c.save()
    buffer.seek(0)
    return buffer.getvalue()


def _add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color="1F2937", align=None):
    if align is None:
        align = PP_ALIGN.LEFT if PP_ALIGN is not None else 0
    box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = box.text_frame
    text_frame.word_wrap = True
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = clean_text(text)
    run.font.size = PptxPt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(color)
    return box


def _add_bullets(slide, left, top, width, height, items, font_size=18, color="1F2937"):
    box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = box.text_frame
    text_frame.word_wrap = True
    text_frame.clear()
    if not items:
        items = ["No details provided."]
    for idx, item in enumerate(items):
        p = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
        p.text = f"• {clean_text(item)}"
        for run in p.runs:
            run.font.size = PptxPt(font_size)
            run.font.color.rgb = RGBColor.from_string(color)
    return box


def _split_lines(value) -> list[str]:
    if isinstance(value, list):
        return [str(item).strip() for item in value if str(item).strip()]
    if isinstance(value, str):
        parts = [part.strip(" •-\t") for part in re.split(r"[\n\r]+", value) if part.strip()]
        return parts or [value.strip()]
    return [str(value).strip()]


def create_lesson_plan_pptx(
    plan_data: dict,
    teacher_name: str = "ISAH YUSUF",
    template_labels: Optional[dict[str, str]] = None,
    template_name: Optional[str] = None,
    theme_key: str = "auto",
) -> bytes:
    if Presentation is None:
        raise HTTPException(status_code=500, detail="PowerPoint generation is unavailable. Install python-pptx.")

    template_labels = template_labels or {}
    theme = resolve_presentation_theme(theme_key, plan_data.get("subject", ""))
    content = dict(plan_data or {})
    prs = Presentation()
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)

    theme_blue = theme["primary"].lstrip("#")
    theme_teal = theme["secondary"].lstrip("#")
    dark = theme["text"].lstrip("#")
    muted = theme["muted"].lstrip("#")
    light_bg = theme["bg"].lstrip("#")

    def add_bg(slide, color="FFFFFF"):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string(color)

    def add_top_bar(slide, title, subtitle=None):
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptxInches(0), PptxInches(0), PptxInches(13.333), PptxInches(0.72))
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor.from_string(theme_blue)
        bar.line.fill.background()
        _add_text_box(slide, PptxInches(0.45), PptxInches(0.12), PptxInches(8.8), PptxInches(0.35), title, font_size=21, bold=True, color="FFFFFF")
        if subtitle:
            _add_text_box(slide, PptxInches(9.05), PptxInches(0.15), PptxInches(3.8), PptxInches(0.25), subtitle, font_size=11, color="DCEBFF", align=PP_ALIGN.RIGHT)

    def add_card(slide, left, top, width, height, fill="F8FAFF", line="D7E3FF"):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor.from_string(fill.lstrip("#"))
        card.line.color.rgb = RGBColor.from_string(line.lstrip("#"))
        return card

    def add_prompt_bar(slide, prompt):
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PptxInches(0.7), PptxInches(6.42), PptxInches(11.95), PptxInches(0.58))
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor.from_string(theme["accent_bg"].lstrip("#"))
        bar.line.color.rgb = RGBColor.from_string(theme["secondary"].lstrip("#"))
        _add_text_box(slide, PptxInches(0.95), PptxInches(6.53), PptxInches(11.4), PptxInches(0.2), f"Teaching prompt: {prompt}", font_size=11, color=theme_blue)

    def add_title_slide():
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slide, light_bg)
        # Decorative brand accents to make the title feel more premium.
        for x, y, w, h, color in [
            (10.65, 0.5, 1.9, 1.9, theme["accent_bg"].lstrip("#")),
            (11.55, 4.95, 1.15, 1.15, theme["secondary"].lstrip("#")),
        ]:
            accent = slide.shapes.add_shape(MSO_SHAPE.OVAL, PptxInches(x), PptxInches(y), PptxInches(w), PptxInches(h))
            accent.fill.solid()
            accent.fill.fore_color.rgb = RGBColor.from_string(color)
            accent.line.color.rgb = RGBColor.from_string(theme["secondary"].lstrip("#"))
        accent = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PptxInches(0.65), PptxInches(0.7), PptxInches(1.2), PptxInches(0.32))
        accent.fill.solid()
        accent.fill.fore_color.rgb = RGBColor.from_string(theme_teal)
        accent.line.fill.background()
        _add_text_box(slide, PptxInches(0.65), PptxInches(1.15), PptxInches(12), PptxInches(0.45), f"{theme['name']} Teaching Deck", font_size=16, bold=True, color=theme_teal)
        _add_text_box(slide, PptxInches(0.65), PptxInches(1.55), PptxInches(12), PptxInches(1.0), f"{content.get('subject', plan_data.get('subject', 'Subject'))} - {content.get('topic', plan_data.get('topic', 'Topic'))}", font_size=28, bold=True, color=dark)
        _add_text_box(slide, PptxInches(0.68), PptxInches(2.48), PptxInches(8.8), PptxInches(0.45), f"Class: {content.get('class', '')}    Date: {content.get('date', '')}", font_size=16, color=muted)
        if template_name:
            _add_text_box(slide, PptxInches(0.68), PptxInches(2.88), PptxInches(8.8), PptxInches(0.35), f"Template: {template_name}", font_size=12, color=theme_teal)
        add_card(slide, PptxInches(9.2), PptxInches(1.05), PptxInches(3.45), PptxInches(3.2), fill="FFFFFF")
        _add_text_box(slide, PptxInches(9.45), PptxInches(1.3), PptxInches(2.95), PptxInches(0.3), "At a glance", font_size=18, bold=True, color=dark, align=PP_ALIGN.CENTER)
        quick_items = [
            f"Topic: {content.get('topic', '')}",
            f"Duration: {content.get('duration', '')}",
            f"Age group: {content.get('age_group', '')}",
            f"Prepared by: {teacher_name}",
        ]
        for idx, line in enumerate(quick_items):
            _add_text_box(slide, PptxInches(9.4), PptxInches(1.78 + idx * 0.46), PptxInches(3.0), PptxInches(0.25), f"• {line}", font_size=12, color=muted)
        add_card(slide, PptxInches(0.65), PptxInches(3.45), PptxInches(12), PptxInches(2.05), fill="FFFFFF")
        _add_text_box(slide, PptxInches(0.95), PptxInches(3.78), PptxInches(11.3), PptxInches(1.3), content.get("cover_subtitle") or content.get("overview_line") or f"Definition, examples, and applications of {content.get('topic', '')}.", font_size=22, color=dark)
        _add_text_box(slide, PptxInches(0.95), PptxInches(5.55), PptxInches(11.3), PptxInches(0.3), theme["hero_note"], font_size=13, color=theme_teal)
        add_prompt_bar(slide, "Open with the topic, then explain why it matters before moving to the meaning.")

    def add_text_slide(title, body, subtitle=None):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slide)
        add_top_bar(slide, title, subtitle)
        add_card(slide, PptxInches(0.7), PptxInches(1.0), PptxInches(11.95), PptxInches(5.95), fill=theme["accent_bg"])
        _add_text_box(slide, PptxInches(1.0), PptxInches(1.35), PptxInches(11.3), PptxInches(0.45), title, font_size=21, bold=True, color=theme_blue)
        _add_text_box(slide, PptxInches(1.0), PptxInches(1.9), PptxInches(11.1), PptxInches(4.2), body, font_size=20, color=dark)
        add_prompt_bar(slide, "Explain the content slowly and check for understanding before moving on.")
        return slide

    def add_bullet_slide(title, items, subtitle=None):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slide)
        add_top_bar(slide, title, subtitle)
        add_card(slide, PptxInches(0.7), PptxInches(1.0), PptxInches(11.95), PptxInches(5.95), fill="FFFFFF")
        _add_text_box(slide, PptxInches(1.0), PptxInches(1.35), PptxInches(11.3), PptxInches(0.4), subtitle or title, font_size=20, bold=True, color=theme_blue)
        _add_bullets(slide, PptxInches(1.0), PptxInches(1.35), PptxInches(11.3), PptxInches(5.3), items, font_size=19, color=dark)
        add_prompt_bar(slide, "Use the content as a quick teaching guide and pause for learners to respond.")
        return slide

    add_title_slide()
    add_text_slide("TOPIC OVERVIEW", content.get("overview_line") or f"{content.get('topic', '')}: Definition, derivation of equations, and applications.", subtitle="Overview")
    add_text_slide(content.get("meaning_heading") or f"MEANING OF {content.get('topic', '').upper()}", content.get("meaning_text") or f"{content.get('topic', '')} is an important concept in {content.get('subject', '')}.", subtitle="Meaning")
    add_bullet_slide(content.get("examples_heading") or "Examples", content.get("examples") or [f"An example related to {content.get('topic', '')}."], subtitle="Examples")
    add_bullet_slide(content.get("key_terms_heading") or "Terms Associated", content.get("key_terms") or [f"Key terms for {content.get('topic', '')}"], subtitle="Terms")
    add_bullet_slide(content.get("worked_examples_heading") or "Worked Examples", content.get("worked_examples") or [f"Worked example for {content.get('topic', '')}."], subtitle="Practice")
    add_bullet_slide(content.get("classwork_heading") or "CLASSWORK", content.get("classwork") or [f"1. Define {content.get('topic', '')}."] , subtitle="Classwork")
    add_bullet_slide(content.get("weekend_assignment_heading") or "WEEKEND ASSIGNMENT", content.get("weekend_assignment") or [f"1. Study {content.get('topic', '')} and answer related questions."], subtitle="Assignment")
    add_text_slide(content.get("closing_line") or "THANK YOU", f"{theme['hero_note']}\n\nPrepared by {teacher_name}", subtitle="Closing")

    bio = io.BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio.getvalue()


@app.on_event("startup")
async def startup_event() -> None:
    init_db()


@app.get("/", response_class=HTMLResponse)
async def home(request: Request):
    return templates.TemplateResponse(
        request,
        "form.html",
        {
            "mode": "lesson",
            "page_title": "Lesson Planner",
            "page_heading": "Generate Lesson Plan",
            "page_subtitle": "Create a formal lesson plan in DOCX or PDF format.",
            "subscription_plans": SUBSCRIPTION_PLANS,
            "presentation_themes": PRESENTATION_THEME_PRESETS,
            "payment_name": PAYMENT_NAME,
            "payment_account": PAYMENT_ACCOUNT,
            "payment_bank": PAYMENT_BANK,
            "admin_path": "/admin",
        },
    )


@app.get("/presentation", response_class=HTMLResponse)
async def presentation_page(request: Request):
    return templates.TemplateResponse(
        request,
        "form.html",
        {
            "mode": "presentation",
            "page_title": "Presentation Builder",
            "page_heading": "Generate Classroom Presentation",
            "page_subtitle": "Create a PPTX slide deck or PDF presentation for teaching and class delivery.",
            "subscription_plans": SUBSCRIPTION_PLANS,
            "presentation_themes": PRESENTATION_THEME_PRESETS,
            "payment_name": PAYMENT_NAME,
            "payment_account": PAYMENT_ACCOUNT,
            "payment_bank": PAYMENT_BANK,
            "admin_path": "/admin",
        },
    )


@app.get("/subscription-history", response_class=HTMLResponse)
async def subscription_history_page(request: Request):
    return templates.TemplateResponse(
        request,
        "subscription_history.html",
        {
            "page_title": "Subscription History",
            "payment_name": PAYMENT_NAME,
            "payment_account": PAYMENT_ACCOUNT,
            "payment_bank": PAYMENT_BANK,
            "admin_path": "/admin",
        },
    )


@app.get("/admin", response_class=HTMLResponse)
async def admin_dashboard(request: Request):
    return templates.TemplateResponse(
        request,
        "admin.html",
        {
            "subscription_plans": SUBSCRIPTION_PLANS,
        },
    )


@app.get("/api/plans")
async def get_plans():
    return {
        "plans": SUBSCRIPTION_PLANS,
        "payment": {"name": PAYMENT_NAME, "account": PAYMENT_ACCOUNT, "bank": PAYMENT_BANK},
    }


@app.get("/api/subscription-status")
async def subscription_status(subscriber_key: str = Query(default="")):
    record = find_subscription(subscriber_key)
    return {"status": record.get("status", "pending") if record else "unregistered", "subscription": record}


@app.get("/api/subscription-history")
async def subscription_history(subscriber_key: str = Query(default="")):
    subscription = find_subscription(subscriber_key)
    if not subscription:
        raise HTTPException(status_code=404, detail="Subscription not found")
    generations = list_generations(subscriber_key, limit=20)
    return {
        "subscription": subscription,
        "generations": generations,
        "summary": {
            "total_generations": len(generations),
            "total_downloads": sum(int(item.get("download_count") or 0) for item in generations),
            "grace_ends_at": subscription.get("grace_ends_at"),
            "days_remaining": subscription.get("days_remaining"),
            "grace_days_remaining": subscription.get("grace_days_remaining"),
        },
    }


@app.post("/api/subscribe")
async def subscribe(payload: SubscriptionRequest):
    record = upsert_subscription(payload)
    return {"message": "Subscription request saved. Awaiting admin approval.", "subscription": record}


@app.get("/api/admin/subscriptions")
async def list_subscriptions(x_admin_password: str = Header(default="", alias="X-Admin-Password")):
    admin_required(x_admin_password)
    return {"items": load_subscriptions()}


@app.get("/api/admin/analytics")
async def admin_analytics(x_admin_password: str = Header(default="", alias="X-Admin-Password")):
    admin_required(x_admin_password)
    return get_generation_analytics()


@app.post("/api/admin/subscriptions/{subscription_id}/approve")
async def approve_subscription(subscription_id: str, x_admin_password: str = Header(default="", alias="X-Admin-Password")):
    admin_required(x_admin_password)
    return {"subscription": set_subscription_status(subscription_id, "approved")}


@app.post("/api/admin/subscriptions/{subscription_id}/reject")
async def reject_subscription(subscription_id: str, x_admin_password: str = Header(default="", alias="X-Admin-Password")):
    admin_required(x_admin_password)
    return {"subscription": set_subscription_status(subscription_id, "rejected")}


@app.get("/api/generations")
async def list_user_generations(subscriber_key: str = Query(default="")):
    subscription = find_subscription(subscriber_key)
    if not subscription or subscription.get("status") != "approved":
        raise HTTPException(status_code=403, detail="Subscription required. Please subscribe before using the app.")
    items = list_generations(subscriber_key)
    return {"items": items}


@app.get("/api/generations/{generation_id}/download")
async def download_generation(generation_id: str, subscriber_key: str = Query(default="")):
    subscription = find_subscription(subscriber_key)
    if not subscription or subscription.get("status") != "approved":
        raise HTTPException(status_code=403, detail="Subscription required. Please subscribe before using the app.")

    generation = get_generation(generation_id)
    if not generation:
        raise HTTPException(status_code=404, detail="Generated file not found")
    if normalize_key(generation.get("subscriber_key")) != normalize_key(subscriber_key):
        raise HTTPException(status_code=403, detail="You do not have access to this file")

    increment_generation_download(generation_id)

    return Response(
        content=generation["file_data"],
        media_type=generation["mime_type"],
        headers={"Content-Disposition": f'attachment; filename="{generation["filename"]}"'},
    )


@app.post("/generate")
async def generate_plan(
    class_name: str = Form(...),
    subject: str = Form(...),
    topic: str = Form(...),
    subscriber_key: str = Form(...),
    output_format: str = Form("docx"),
    content_type: str = Form("lesson"),
    presentation_theme: str = Form("auto"),
    lesson_template: UploadFile | None = File(None),
):
    if not subscriber_key:
        raise HTTPException(status_code=403, detail="Subscription required. Please subscribe before using the app.")

    subscription = find_subscription(subscriber_key)
    if not subscription:
        raise HTTPException(status_code=403, detail="Subscription required. Please subscribe before using the app.")
    if subscription.get("status") == "expired":
        renewal = subscription.get("next_renewal_at") or subscription.get("expires_at") or "the renewal date"
        raise HTTPException(status_code=403, detail=f"Subscription expired. Please resubscribe to renew access from {renewal}.")
    if subscription.get("status") != "approved":
        raise HTTPException(status_code=403, detail="Subscription pending. Access is unlocked after admin approval.")

    template_text = ""
    template_name = None
    template_labels = {}
    if lesson_template and lesson_template.filename:
        template_text, _ = extract_template_outline(lesson_template)
        template_outline = derive_template_outline(template_text)
        template_name = lesson_template.filename
        template_labels = template_outline["labels"] if template_text.strip() else {}
        template_prompt_text = (
            f"Detected section order: {', '.join(template_outline['sections']) or 'default lesson structure'}\n\n"
            f"Template text:\n{template_outline['prompt_text']}"
        )
    else:
        template_prompt_text = ""

    is_presentation = normalize_key(content_type) == "presentation"
    if is_presentation:
        plan_data = ai_service.generate_presentation_content(
            subject=subject,
            class_level=class_name,
            topic=topic,
            template_outline=template_prompt_text,
        )
    else:
        plan_data = ai_service.generate_lesson_plan(
            subject=subject,
            class_level=class_name,
            topic=topic,
            template_outline=template_prompt_text,
        )

    plan_data.setdefault("class", class_name)
    plan_data.setdefault("subject", subject)
    plan_data.setdefault("topic", topic)
    plan_data.setdefault("date", datetime.now().strftime("%d %B, %Y"))
    plan_data.setdefault("week", "Two")
    plan_data.setdefault("duration", "Forty Minutes")
    plan_data.setdefault("age_group", f"{class_name} students")
    if is_presentation:
        plan_data.setdefault("cover_subtitle", f"Definition, examples, and applications of {topic}.")
        plan_data.setdefault("overview_line", f"{topic}: Definition, derivation of equations, and applications.")
        plan_data.setdefault("meaning_heading", f"MEANING OF {topic.upper()}")
        plan_data.setdefault("meaning_text", f"{topic} is an important concept in {subject} for {class_name} learners.")
        plan_data.setdefault("examples_heading", "Examples")
        plan_data.setdefault("examples", [f"An example related to {topic}."])
        plan_data.setdefault("key_terms_heading", f"TERMS ASSOCIATED WITH {topic.upper()}")
        plan_data.setdefault("key_terms", [f"Key term one in {topic}"])
        plan_data.setdefault("worked_examples_heading", "Examples")
        plan_data.setdefault("worked_examples", [f"Worked example on {topic}."])
        plan_data.setdefault("classwork_heading", "CLASSWORK")
        plan_data.setdefault("classwork", [f"1. Define {topic} and mention two applications."])
        plan_data.setdefault("weekend_assignment_heading", "WEEKEND ASSIGNMENT")
        plan_data.setdefault("weekend_assignment", [f"1. Study {topic} and answer related questions."])
        plan_data.setdefault("closing_line", "THANK YOU")

    safe_subject = sanitize_filename(subject, max_len=30)
    safe_topic = sanitize_filename(topic, max_len=40)
    selected_format = normalize_output_format(output_format)
    content_type = normalize_key(content_type)

    if selected_format == "pptx":
        file_bytes = create_lesson_plan_pptx(
            plan_data,
            teacher_name="ISAH YUSUF",
            template_labels=template_labels,
            template_name=template_name,
            theme_key=presentation_theme,
        )
        spec = get_output_spec(selected_format)
        filename = f"{spec['prefix']}_{safe_subject}_{safe_topic}.{spec['extension']}"
        mime_type = spec["mime_type"]
    elif selected_format == "pdf" and content_type == "presentation":
        file_bytes = create_presentation_pdf(
            plan_data,
            teacher_name="ISAH YUSUF",
            template_labels=template_labels,
            template_name=template_name,
            theme_key=presentation_theme,
        )
        spec = get_output_spec(selected_format)
        filename = f"Presentation_{safe_subject}_{safe_topic}.{spec['extension']}"
        mime_type = spec["mime_type"]
    elif selected_format == "pdf":
        file_bytes = create_lesson_plan_pdf(
            plan_data,
            teacher_name="ISAH YUSUF",
            template_labels=template_labels,
            template_name=template_name,
        )
        spec = get_output_spec(selected_format)
        filename = f"Lesson_Export_{safe_subject}_{safe_topic}.{spec['extension']}"
        mime_type = spec["mime_type"]
    else:
        file_bytes = create_lesson_plan_doc(
            plan_data,
            teacher_name="ISAH YUSUF",
            template_labels=template_labels,
            template_name=template_name,
        )
        spec = get_output_spec(selected_format)
        filename = f"{spec['prefix']}_{safe_subject}_{safe_topic}.{spec['extension']}"
        mime_type = spec["mime_type"]

    save_generation(
        subscriber_key=subscriber_key,
        output_format=selected_format,
        presentation_theme=presentation_theme,
        subject=subject,
        topic=topic,
        filename=filename,
        mime_type=mime_type,
        file_bytes=file_bytes,
    )
    return Response(
        content=file_bytes,
        media_type=mime_type,
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )
